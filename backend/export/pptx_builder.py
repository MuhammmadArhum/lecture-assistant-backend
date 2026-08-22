"""
Builds a themed .pptx from a Final Brief JSON object.

Design system: card-panel content blocks, numbered chips, and small tag
labels instead of plain bullet lists, with the actual color palette pulled
from `backend.themes` (a handful of presets -- Midnight/Daylight/Sunset/
Ocean/Forest/Slate) so the user can pick a look the way Gamma lets you pick
a theme before generating. Kept isolated from main.py per the assignment's
module-separation requirement. Every key-finding citation is still carried
through onto its slide.

`build_slide_plan()` re-uses the same pure-python content-assembly helpers
(_segment_content_units / _build_segment_slides / _chunk) to hand back a
JSON-serializable slide-by-slide plan with no python-pptx objects in it --
this is what powers the frontend's in-browser deck preview, and it is kept
in lockstep with `build_pptx()` by construction (same helpers, same order).
"""
from __future__ import annotations

import io
import re
import threading
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls

from backend.themes import DEFAULT_THEME, get_theme

# --- Palette -----------------------------------------------------------
# These module-level colors are the *active* theme's palette. They start
# out as the default theme and are reassigned by `_apply_theme()` right
# before each build_pptx()/build_slide_plan() call -- every place below
# that needs a color reads one of these names directly (not as a stale
# function-default value; see the `color=None` + in-body resolution
# pattern used by _text/_accent_bar/_title_box/_header/_text_slide) so a
# theme switch actually takes effect. A module-level lock serializes
# builds so two concurrent requests with different themes can't race.
_BUILD_LOCK = threading.Lock()


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = (hex_str or "#000000").lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_theme(theme_id: str | None) -> dict:
    theme = get_theme(theme_id)
    g = globals()
    g["BACKGROUND"] = _hex_to_rgb(theme["background"])
    g["CARD_BG"] = _hex_to_rgb(theme["card_bg"])
    g["CARD_BG_ALT"] = _hex_to_rgb(theme["card_bg_alt"])
    g["INK"] = _hex_to_rgb(theme["ink"])
    g["MUTED"] = _hex_to_rgb(theme["muted"])
    g["CHIP_TEXT"] = _hex_to_rgb(theme["chip_text"])
    accents = [_hex_to_rgb(c) for c in theme["accents"]]
    g["ACCENT"], g["ACCENT2"], g["ACCENT3"] = accents[0], accents[1], accents[2]
    g["_ACCENTS"] = accents
    g["TITLE_FONT"] = theme.get("title_font") or g.get("TITLE_FONT", "Segoe UI Semibold")
    g["LABEL_FONT"] = theme.get("label_font") or g.get("LABEL_FONT", "Consolas")
    g["BODY_FONT"] = theme.get("body_font") or g.get("BODY_FONT", "Segoe UI")
    return theme


def _accent_for(i: int) -> RGBColor:
    return _ACCENTS[i % len(_ACCENTS)]


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Seed the module-level palette/fonts with the default theme so anything
# imported before the first real build_pptx() call still has valid colors.
_apply_theme(DEFAULT_THEME)

# --- Dynamic slide count -------------------------------------------------
# Calibrated against two instructor-given reference points:
#   10 minutes -> ~15 slides,  60 minutes -> 65-75 slides (~70 midpoint)
# Solving total = a*minutes + b for those two points gives a=1.1, b=4.
SLIDES_PER_MINUTE = 1.1
FIXED_SLIDE_OVERHEAD = 4
MIN_TOTAL_SLIDES = 6


def _target_slide_count(target_minutes: int | None) -> int:
    minutes = target_minutes or 45
    return max(MIN_TOTAL_SLIDES, round(minutes * SLIDES_PER_MINUTE) + FIXED_SLIDE_OVERHEAD)


def _chunk(items: list, size: int) -> list[list]:
    if size < 1:
        size = 1
    return [items[i:i + size] for i in range(0, len(items), size)] or [[]]


# --- Text sanitation -------------------------------------------------------
# Hard cap on how much text goes on one slide/card, and the font size used
# for it. Content should already be a few clean sentences (the prompts ask
# for that), but this is a last line of defense against any upstream node
# handing over a long, unprocessed block of text -- keeps a slide readable
# and on-slide instead of running off the edge.
_BODY_MAX_CHARS = 620


def _clean_slide_text(raw: str, max_chars: int = _BODY_MAX_CHARS) -> str:
    text = raw or ""
    text = re.sub(r"\[\.\.\.\]", " ", text)
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    # Drop byline lines ("By Cole Stryker, Mark Scapicchio") and bare
    # nav/footer lines ("Twitter/X", "Contact Us") that survive scraping --
    # these are page furniture, not lecture content.
    text = re.sub(r"^By\s+[A-Z][\w.'-]*(?:\s+[A-Z][\w.'-]*)*(,.*)?\s*$", "", text, flags=re.MULTILINE)
    text = "\n".join(
        line for line in text.splitlines()
        if not re.fullmatch(r"(Twitter/X|YouTube|LinkedIn|Facebook|Instagram|Contact Us)", line.strip())
    )
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if len(text) <= max_chars:
        return text
    truncated = text[:max_chars]
    last_period = truncated.rfind(". ")
    if last_period > max_chars * 0.4:
        return truncated[:last_period + 1]
    return truncated.rsplit(" ", 1)[0] + "..."


def _body_font_size(text: str, base: int = 19) -> int:
    # Scale down as content approaches the cap so long-but-valid content
    # still fits comfortably rather than relying only on PowerPoint's
    # auto-shrink at open time.
    length = len(text)
    if length > 460:
        return max(13, base - 4)
    if length > 260:
        return max(14, base - 2)
    return base


# --- Low-level shape helpers -------------------------------------------------

def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BACKGROUND
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    _add_slide_transition(slide)
    return slide


# --- Animation (raw-XML) -----------------------------------------------------
# python-pptx has no animation API at all -- PowerPoint animations live in a
# <p:timing> block that has to be hand-built as OOXML and spliced onto each
# slide's XML tree. This section is the only place that does that.

def _add_slide_transition(slide, effect: str = "fade", speed: str = "med"):
    """Adds a between-slide transition (plays when advancing to this slide)."""
    sld = slide._element
    existing = sld.find(qn("p:transition"))
    if existing is not None:
        sld.remove(existing)
    transition_el = parse_xml(f'<p:transition {nsdecls("p")} spd="{speed}"><p:fade/></p:transition>')
    timing = sld.find(qn("p:timing"))
    if timing is not None:
        timing.addprevious(transition_el)
    else:
        sld.append(transition_el)


def _apply_build_animations(slide, stages: list[list]):
    """Adds click-triggered "Fade In" build animations to slide content.

    `stages` is an ordered list of shape-groups. Each group's shapes fade in
    together (as one visual unit -- e.g. a whole card); each new group
    requires an additional click to reveal, in order. Groups/shapes with no
    content are skipped automatically. No-op if there's nothing to animate.
    """
    stages = [[shp for shp in group if shp is not None] for group in stages]
    stages = [group for group in stages if group]
    if not stages:
        return

    next_id = [10]

    def _id():
        next_id[0] += 1
        return next_id[0]

    par_nodes, bld_nodes = [], []
    for group in stages:
        for j, shape in enumerate(group):
            spid = shape.shape_id
            node_type = "clickEffect" if j == 0 else "withEffect"
            par_nodes.append(f'''
            <p:par>
              <p:cTn id="{_id()}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{node_type}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{_id()}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animEffect transition="in" filter="fade">
                    <p:cBhvr>
                      <p:cTn id="{_id()}" dur="500"/>
                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>''')
            bld_nodes.append(f'<p:bldP spid="{spid}" grpId="0"/>')

    timing_xml = f'''<p:timing {nsdecls("p")}>
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    {''.join(par_nodes)}
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
      <p:bldLst>
        {''.join(bld_nodes)}
      </p:bldLst>
    </p:timing>'''

    sld = slide._element
    existing = sld.find(qn("p:timing"))
    if existing is not None:
        sld.remove(existing)
    sld.append(parse_xml(timing_xml))


def _rect(slide, left, top, width, height, fill: RGBColor, rounded=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _oval(slide, left, top, diameter, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(box_or_tf, text, *, size=18, color=None, bold=False, italic=False,
          font=None, align=None, first=True, space_after=6,
          hyperlink: str | None = None, underline: bool = False):
    # color/font default to the *current* theme's INK/BODY_FONT -- resolved
    # here (a live global lookup) rather than as literal default-argument
    # values, which would freeze onto whatever theme was active at import
    # time and never update when _apply_theme() switches themes.
    if color is None:
        color = INK
    if font is None:
        font = BODY_FONT
    tf = box_or_tf if hasattr(box_or_tf, "add_paragraph") else box_or_tf.text_frame
    p = tf.paragraphs[0] if (first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    if hyperlink:
        # Makes the run an actual clickable link in the exported .pptx
        # (opens in the default browser when clicked in PowerPoint).
        run.hyperlink.address = hyperlink
        run.font.underline = underline
    elif underline:
        run.font.underline = True
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    return p


def _normalize_url(raw: str) -> str:
    url = (raw or "").strip()
    if url and not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", url):
        url = "https://" + url
    return url


def _url_from_citation(citation: str) -> str | None:
    """Citations are formatted 'Source Title (url)' -- pull the url back out
    so it can be made clickable without needing a separate stored field."""
    match = re.search(r"\(([^()]*)\)\s*$", citation or "")
    if not match:
        return None
    candidate = match.group(1).strip()
    if "." not in candidate or " " in candidate:
        return None
    return _normalize_url(candidate)


def _tag_chip(slide, text: str, left=Inches(0.7), top=Inches(0.5), color=None, dark_bg=False):
    """Small pill/rect chip with solid accent fill and bold uppercase label
    -- mirrors the reference deck's 'CASE STUDY' tag."""
    color = color or ACCENT
    width = Inches(0.35 + 0.115 * len(text))
    height = Inches(0.34)
    chip = _rect(slide, left, top, width, height, fill=(CARD_BG_ALT if dark_bg else color), rounded=True, radius=0.5)
    tf = chip.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _text(tf, text.upper(), size=11, color=(color if dark_bg else CHIP_TEXT),
          bold=True, font=LABEL_FONT, align=PP_ALIGN.CENTER, space_after=0)
    return chip, width


def _numbered_chip(slide, label: str, left, top, color, diameter=Inches(0.5)):
    """Small circular chip with a bold number/glyph -- used as a per-slide
    or per-card index marker, rotating through the accent palette."""
    chip = _oval(slide, left, top, diameter, fill=color)
    tf = chip.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _text(tf, label, size=16, color=CHIP_TEXT, bold=True, font=LABEL_FONT, align=PP_ALIGN.CENTER, space_after=0)
    return chip


def _accent_bar(slide, left, top, width=Inches(0.9), height=Inches(0.07), color=None):
    if color is None:
        color = ACCENT
    return _rect(slide, left, top, width, height, fill=color)


def _title_box(slide, text: str, top=Inches(0.95), size=30, color=None, width=Inches(11.9), left=Inches(0.7)):
    if color is None:
        color = INK
    box = slide.shapes.add_textbox(left, top, width, Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = TITLE_FONT
    return box


def _header(slide, eyebrow: str, title: str, accent: RGBColor = None, title_size=30):
    """Standard slide header: tag chip + title + thin accent bar underneath."""
    if accent is None:
        accent = ACCENT
    _tag_chip(slide, eyebrow, left=Inches(0.7), top=Inches(0.5), color=accent)
    _title_box(slide, title, top=Inches(1.0), size=title_size)
    _accent_bar(slide, Inches(0.7), Inches(1.82), width=Inches(0.9), color=accent)


def _decorative_corner(slide):
    """Small tasteful geometric flourish, bottom-right corner -- a trio of
    accent dots at falling sizes, echoing the reference deck's icon/accent
    treatment without needing external images."""
    base_x, base_y = Inches(12.55), Inches(6.85)
    sizes = [Inches(0.14), Inches(0.09), Inches(0.055)]
    for i, size in enumerate(sizes):
        _oval(slide, base_x + Inches(0.22 * i), base_y, size, fill=_accent_for(i))


# --- Slide builders ----------------------------------------------------------

def _title_slide(prs: Presentation, brief: dict[str, Any]):
    slide = _blank_slide(prs)

    # Large soft accent circle bleeding off the top-right edge for visual
    # weight, echoed by a smaller one lower-left -- gives the title slide
    # some presence without needing an external image.
    _oval(slide, Inches(10.6), Inches(-1.6), Inches(4.2), fill=CARD_BG_ALT)
    _oval(slide, Inches(-1.2), Inches(5.6), Inches(3.0), fill=CARD_BG)

    _tag_chip(slide, "LECTURE BRIEF", left=Inches(0.9), top=Inches(1.5), color=ACCENT)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(11.3), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = brief.get("title", "Untitled Lecture")
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = INK
    p.font.name = TITLE_FONT
    p.alignment = PP_ALIGN.LEFT

    _accent_bar(slide, Inches(0.9), Inches(4.05), width=Inches(1.5), color=ACCENT)

    minutes = brief.get("target_minutes")
    if minutes:
        _text(
            slide.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(9), Inches(0.5)),
            f"{minutes}-minute session",
            size=15, color=MUTED, font=BODY_FONT,
        )


def _text_slide(prs: Presentation, eyebrow: str, title: str, body: str, accent: RGBColor = None):
    if accent is None:
        accent = ACCENT
    slide = _blank_slide(prs)
    _header(slide, eyebrow, title, accent=accent)

    cleaned = _clean_slide_text(body, max_chars=1000)
    card = _rect(slide, Inches(0.7), Inches(2.15), Inches(11.9), Inches(4.55), fill=CARD_BG, rounded=True, radius=0.045)
    card.text_frame.word_wrap = True
    card.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    card.text_frame.margin_left = Inches(0.5)
    card.text_frame.margin_right = Inches(0.5)
    card.text_frame.margin_top = Inches(0.45)
    _text(card.text_frame, cleaned, size=_body_font_size(cleaned, base=20), color=MUTED, font=BODY_FONT, space_after=0)
    _decorative_corner(slide)
    _apply_build_animations(slide, [[card]])


def _focus_slide(prs: Presentation, eyebrow: str, index: int, title: str, body: str, accent: RGBColor):
    """A segment slide with exactly one subtopic -- single big focus card."""
    slide = _blank_slide(prs)
    _tag_chip(slide, eyebrow, left=Inches(0.7), top=Inches(0.5), color=accent)
    _numbered_chip(slide, str(index), left=Inches(11.9), top=Inches(0.45), color=accent, diameter=Inches(0.45))
    _title_box(slide, title, top=Inches(1.0), size=28, width=Inches(11.0))
    _accent_bar(slide, Inches(0.7), Inches(1.95), width=Inches(0.9), color=accent)

    cleaned = _clean_slide_text(body)
    card = _rect(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.4), fill=CARD_BG, rounded=True, radius=0.05)
    card.text_frame.word_wrap = True
    card.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    card.text_frame.margin_left = Inches(0.5)
    card.text_frame.margin_right = Inches(0.5)
    card.text_frame.margin_top = Inches(0.4)
    _text(card.text_frame, cleaned, size=_body_font_size(cleaned, base=20), color=MUTED, font=BODY_FONT, space_after=0)
    _apply_build_animations(slide, [[card]])


def _card_grid_slide(prs: Presentation, eyebrow: str, title: str, cards: list[dict[str, str]],
                      columns: int = 2, chip_glyphs: list[str] | None = None):
    """Grid of rounded-rect cards, each with a numbered/glyph chip, a bold
    heading, a body, and an optional muted footer (citation/url). Used for
    multi-subtopic segment slides, key findings, risks, and further reading
    -- this is the reference deck's card-grid language generalized to any
    list of items."""
    slide = _blank_slide(prs)
    _header(slide, eyebrow, title)

    n = len(cards)
    columns = max(1, min(columns, n))
    rows = -(-n // columns)

    area_left, area_top = Inches(0.7), Inches(2.15)
    area_w, area_h = Inches(11.9), Inches(4.55)
    gap = Inches(0.28)
    card_w = Emu(int((area_w - gap * (columns - 1)) / columns))
    card_h = Emu(int((area_h - gap * (rows - 1)) / rows))

    body_cap = 340 if n <= 2 else (220 if n <= 4 else 150)

    card_stages: list[list] = []

    for i, c in enumerate(cards):
        col, row = i % columns, i // columns
        left = area_left + col * (card_w + gap)
        top = area_top + row * (card_h + gap)
        accent = _accent_for(i)

        card_rect = _rect(slide, left, top, card_w, card_h, fill=CARD_BG, rounded=True, radius=0.06)
        card_shapes = [card_rect]

        glyph = chip_glyphs[i] if chip_glyphs else str(i + 1)
        chip_d = Inches(0.42)
        chip = _numbered_chip(slide, glyph, left + Inches(0.28), top + Inches(0.25), color=accent, diameter=chip_d)
        card_shapes.append(chip)

        heading = _clean_slide_text(c.get("heading", ""), max_chars=70)
        head_box = slide.shapes.add_textbox(left + Inches(0.85), top + Inches(0.22), card_w - Inches(1.1), Inches(0.5))
        htf = head_box.text_frame
        htf.word_wrap = True
        _text(htf, heading, size=16, color=INK, bold=True, font=TITLE_FONT, space_after=0)
        card_shapes.append(head_box)

        footer = c.get("footer")
        footer_text = _clean_slide_text(footer, max_chars=90) if footer else ""
        foot_h = (Inches(0.55) if len(footer_text) > 45 else Inches(0.32)) if footer_text else Inches(0)
        foot_reserve = (foot_h + Inches(0.1)) if footer_text else Inches(0)

        body_text = _clean_slide_text(c.get("body", ""), max_chars=body_cap)
        body_box = slide.shapes.add_textbox(
            left + Inches(0.28), top + Inches(0.85),
            card_w - Inches(0.56), card_h - Inches(1.05) - foot_reserve,
        )
        btf = body_box.text_frame
        btf.word_wrap = True
        btf.auto_size = MSO_AUTO_SIZE.NONE
        # If this card is pointing at a URL (further-reading cards pass one
        # explicitly; others may have one embedded in the body text itself),
        # make the body text an actual clickable link.
        body_url = c.get("url") or (_normalize_url(body_text) if re.match(r"^(https?://|www\.)", body_text) else None)
        _text(btf, body_text, size=_body_font_size(body_text, base=14),
              color=(accent if body_url else MUTED), font=BODY_FONT, space_after=0,
              hyperlink=body_url)
        card_shapes.append(body_box)

        if footer_text:
            footer_url = c.get("footer_url") or _url_from_citation(footer)
            foot_box = slide.shapes.add_textbox(
                left + Inches(0.28), top + card_h - foot_h - Inches(0.08),
                card_w - Inches(0.56), foot_h,
            )
            ftf = foot_box.text_frame
            ftf.word_wrap = True
            ftf.auto_size = MSO_AUTO_SIZE.NONE
            _text(ftf, footer_text, size=10, color=accent,
                  italic=True, font=BODY_FONT, space_after=0, hyperlink=footer_url)
            card_shapes.append(foot_box)

        card_stages.append(card_shapes)

    _apply_build_animations(slide, card_stages)


def _appendix_slide(prs: Presentation, node_trace: list[dict[str, Any]]):
    slide = _blank_slide(prs)
    _header(slide, "Appendix", "Node Trace", accent=ACCENT3)
    card = _rect(slide, Inches(0.7), Inches(2.15), Inches(11.9), Inches(4.55), fill=CARD_BG, rounded=True, radius=0.045)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.35)
    for i, entry in enumerate(node_trace[:12]):
        node = entry.get("node", "?")
        ts = entry.get("timestamp", "")
        decision = entry.get("human_decision")
        suffix = f"   |   human: {decision}" if decision else ""
        _text(tf, f"{ts}  —  {node}{suffix}", size=12, color=MUTED, font=LABEL_FONT, space_after=8, first=(i == 0))
    _apply_build_animations(slide, [[card]])


# --- Content assembly (unchanged logic from prior fix, re-skinned above) ---

def _segment_content_units(segment: dict[str, Any]) -> list[dict[str, str]]:
    """Real, presentable content units for one segment, each destined for its
    own slide/card. Prefers the LLM-authored `subtopics` (title + real
    content per slide). Falls back to splitting the segment's short `notes`
    paragraph by sentence for older plans that predate subtopics -- but
    never manufactures empty units; a segment simply contributes as many
    real slides as it has real content for."""
    subtopics = segment.get("subtopics") or []
    units = [
        {"title": s.get("title") or segment.get("label", ""), "content": (s.get("content") or "").strip()}
        for s in subtopics
        if (s.get("content") or "").strip()
    ]
    if units:
        return units

    notes = (segment.get("notes") or "").strip()
    if not notes:
        return [{"title": segment.get("label", ""), "content": segment.get("label", "")}]
    sentences = [p.strip() for p in notes.replace("!", ".").split(".") if p.strip()]
    return [{"title": segment.get("label", ""), "content": s + "."} for s in sentences] or \
        [{"title": segment.get("label", ""), "content": notes}]


def _build_segment_slides(segments: list[dict[str, Any]], slide_budget: int) -> list[dict[str, Any]]:
    """Turn N plan segments into up to slide_budget slide-entries. Each
    entry carries its raw `units` (1 or more {title, content} dicts) rather
    than pre-joined text, so the renderer can choose a single-focus layout
    for one unit or a card grid for several. Distributes the budget across
    segments proportional to each segment's share of minutes, but a segment
    is never padded past the amount of real content it actually has -- any
    unused budget is handed to other segments (via grouping their extra
    subtopics onto shared slides) rather than emitted as blank slides."""
    if not segments or slide_budget <= 0:
        return []

    total_minutes = sum(s.get("minutes") or 0 for s in segments) or len(segments)
    content = [_segment_content_units(s) for s in segments]

    counts = []
    for s, units in zip(segments, content):
        share = (s.get("minutes") or (total_minutes / len(segments))) / total_minutes
        counts.append(min(len(units), max(1, round(slide_budget * share))))

    total_assigned = sum(counts)
    if total_assigned > slide_budget:
        while total_assigned > slide_budget and any(c > 1 for c in counts):
            i = counts.index(max(counts))
            counts[i] -= 1
            total_assigned -= 1

    entries = []
    for seg, units, count in zip(segments, content, counts):
        if count <= 0:
            continue
        groups = _chunk(units, max(1, -(-len(units) // count)))[:count]  # ceil-based grouping, no padding
        for group in groups:
            entries.append({"segment": seg, "units": group})
    return entries


def build_pptx(brief: dict[str, Any], theme: str | None = None) -> io.BytesIO:
    with _BUILD_LOCK:
        return _build_pptx_locked(brief, theme)


def _build_pptx_locked(brief: dict[str, Any], theme: str | None) -> io.BytesIO:
    _apply_theme(theme or brief.get("theme"))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    target_total = _target_slide_count(brief.get("target_minutes"))

    introduction = (brief.get("introduction") or "").strip()
    summary = brief.get("summary", "")
    segments = brief.get("segments") or []
    findings = brief.get("key_findings", [])
    risks = brief.get("risks", [])
    reading = brief.get("further_reading", [])
    node_trace = brief.get("node_trace")

    fixed_front = 1 + (1 if introduction else 0) + 1
    # Findings/reading grid comfortably at up to 4 per slide (2x2); risks
    # (usually short, punchy) also 4 per slide.
    findings_chunks = _chunk(findings, size=4) if findings else []
    risk_chunks = _chunk(risks, size=4) if risks else []
    reading_chunks = _chunk(reading, size=4) if reading else []
    appendix_count = 1 if node_trace else 0
    tail_count = len(findings_chunks) + len(risk_chunks) + len(reading_chunks) + appendix_count

    segment_budget = max(0, target_total - fixed_front - tail_count)
    segment_entries = _build_segment_slides(segments, segment_budget)

    # --- Emit slides in order ---
    _title_slide(prs, brief)
    if introduction:
        _text_slide(prs, "Part 1", "Introduction", introduction, accent=ACCENT)
    _text_slide(prs, "Part 2", "Summary", summary, accent=ACCENT2)

    n_segment_slides = len(segment_entries)
    for i, entry in enumerate(segment_entries, start=1):
        seg = entry["segment"]
        units = entry["units"]
        eyebrow = f"Part 3 · {i}/{n_segment_slides}"
        accent = _accent_for(i - 1)
        if len(units) == 1:
            _focus_slide(prs, eyebrow, i, units[0]["title"] or seg.get("label", ""), units[0]["content"], accent)
        else:
            cards = [{"heading": u["title"], "body": u["content"]} for u in units]
            _card_grid_slide(prs, eyebrow, seg.get("label", "Untitled segment"), cards, columns=2)

    for chunk in findings_chunks:
        cards = [
            {
                "heading": (f.get("citation", "").split(" (")[0].strip() or "Finding"),
                "body": f.get("text", ""),
                "footer": f.get("citation", ""),
            }
            for f in chunk
        ]
        _card_grid_slide(prs, "Part 4", "Key Findings", cards, columns=2)
    for chunk in risk_chunks:
        cards = [{"heading": "Risk", "body": r} for r in chunk]
        _card_grid_slide(prs, "Part 5", "Risks", cards, columns=2, chip_glyphs=["!"] * len(chunk))
    for chunk in reading_chunks:
        cards = [{"heading": item.get("title", ""), "body": item.get("url", ""), "url": _normalize_url(item.get("url", "")), "footer": None} for item in chunk]
        _card_grid_slide(prs, "Part 6", "Further Reading", cards, columns=2, chip_glyphs=["→"] * len(chunk))

    if node_trace:
        _appendix_slide(prs, node_trace)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def build_slide_plan(brief: dict[str, Any]) -> list[dict[str, Any]]:
    """Pure-data mirror of build_pptx()'s slide assembly (no python-pptx
    objects), for the frontend's in-browser deck preview. Shares the same
    _segment_content_units/_build_segment_slides/_chunk helpers and the
    same ordering so the preview never drifts from the real export.

    Each slide dict has a "kind" of "title" | "text" | "focus" | "cards" |
    "appendix", plus an "accent_index" (0/1/2) the frontend maps onto the
    chosen theme's 3 accent colors, matching the rotation _accent_for()
    uses in the real render.
    """
    target_total = _target_slide_count(brief.get("target_minutes"))

    introduction = (brief.get("introduction") or "").strip()
    summary = brief.get("summary", "")
    segments = brief.get("segments") or []
    findings = brief.get("key_findings", [])
    risks = brief.get("risks", [])
    reading = brief.get("further_reading", [])
    node_trace = brief.get("node_trace")

    fixed_front = 1 + (1 if introduction else 0) + 1
    findings_chunks = _chunk(findings, size=4) if findings else []
    risk_chunks = _chunk(risks, size=4) if risks else []
    reading_chunks = _chunk(reading, size=4) if reading else []
    appendix_count = 1 if node_trace else 0
    tail_count = len(findings_chunks) + len(risk_chunks) + len(reading_chunks) + appendix_count

    segment_budget = max(0, target_total - fixed_front - tail_count)
    segment_entries = _build_segment_slides(segments, segment_budget)

    slides: list[dict[str, Any]] = []

    slides.append({
        "kind": "title",
        "eyebrow": "Lecture Brief",
        "title": brief.get("title", "Untitled Lecture"),
        "subtitle": f'{brief.get("target_minutes")}-minute session' if brief.get("target_minutes") else None,
    })

    if introduction:
        slides.append({
            "kind": "text", "eyebrow": "Part 1", "title": "Introduction",
            "body": _clean_slide_text(introduction, max_chars=1000), "accent_index": 0,
        })

    slides.append({
        "kind": "text", "eyebrow": "Part 2", "title": "Summary",
        "body": _clean_slide_text(summary or "", max_chars=1000), "accent_index": 1,
    })

    n_segment_slides = len(segment_entries)
    for i, entry in enumerate(segment_entries, start=1):
        seg = entry["segment"]
        units = entry["units"]
        eyebrow = f"Part 3 · {i}/{n_segment_slides}"
        accent_index = (i - 1) % 3
        if len(units) == 1:
            slides.append({
                "kind": "focus", "eyebrow": eyebrow, "index": i,
                "title": units[0]["title"] or seg.get("label", ""),
                "body": _clean_slide_text(units[0]["content"]),
                "accent_index": accent_index,
            })
        else:
            cards = [
                {"heading": _clean_slide_text(u["title"], max_chars=70),
                 "body": _clean_slide_text(u["content"], max_chars=220)}
                for u in units
            ]
            slides.append({
                "kind": "cards", "eyebrow": eyebrow,
                "title": seg.get("label", "Untitled segment"), "cards": cards, "columns": 2,
            })

    for chunk in findings_chunks:
        cards = [
            {
                "heading": (f.get("citation", "").split(" (")[0].strip() or "Finding"),
                "body": _clean_slide_text(f.get("text", ""), max_chars=220),
                "footer": f.get("citation", ""),
            }
            for f in chunk
        ]
        slides.append({"kind": "cards", "eyebrow": "Part 4", "title": "Key Findings", "cards": cards, "columns": 2})

    for chunk in risk_chunks:
        cards = [{"heading": "Risk", "body": _clean_slide_text(r, max_chars=220)} for r in chunk]
        slides.append({"kind": "cards", "eyebrow": "Part 5", "title": "Risks", "cards": cards, "columns": 2})

    for chunk in reading_chunks:
        cards = [
            {"heading": item.get("title", ""), "body": item.get("url", ""), "url": _normalize_url(item.get("url", ""))}
            for item in chunk
        ]
        slides.append({"kind": "cards", "eyebrow": "Part 6", "title": "Further Reading", "cards": cards, "columns": 2})

    if node_trace:
        slides.append({"kind": "appendix", "eyebrow": "Appendix", "title": "Node Trace", "trace": node_trace[:12]})

    return slides
